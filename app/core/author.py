# -*- coding: utf-8 -*-
from __future__ import annotations
import os, re, json, time
from pathlib import Path
from typing import List, Dict, Any, Optional
from typing import List, Dict, Any
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from typing import List, Dict

# Gemini クライアント（ある場合だけ使う）
try:
    from app.core.gemini_client import get_client  # get_client() -> google.genai.client.Client
except Exception:
    get_client = None  # 型だけ合わせておく


def _normalize_bullets(bullets: List[str], max_chars: int = 40) -> List[str]:
    """
    箇条書きをスライド用に整形する:
      - 前後の空白除去
      - 文末の句読点(。．.)を削る
      - 長すぎる場合は max_chars でカットして "…" を付ける
    """
    normalized: List[str] = []
    for b in bullets:
        s = str(b).strip()
        if not s:
            continue
        # 文末の句読点などを削る
        s = re.sub(r"[。、．\.]+$", "", s)
        # 長すぎる行はカット
        if len(s) > max_chars:
            s = s[:max_chars] + "…"
        normalized.append(s)
    return normalized



# ========== ユーティリティ ==========

def _mk_pres() -> Presentation:
    return Presentation()

def _safe_save(prs: Presentation, out_path: str | Path) -> str:
    p = Path(out_path)
    p.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(p))
    return str(p)

def _add_title_slide(prs: Presentation, title: str, subtitle: Optional[str] = None):
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title
    slide.shapes.title.text = title
    sub = slide.placeholders[1]
    if not subtitle:
        subtitle = time.strftime("%Y-%m-%d %H:%M")
    sub.text = subtitle
    return slide

def _add_bullet_slide(prs: Presentation, title: str, bullets: List[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    if not bullets:
        bullets = ["（内容なし）"]
    # 最初の段落
    tf.text = str(bullets[0]).strip()
    # 2行目以降
    for b in bullets[1:]:
        p = tf.add_paragraph()
        p.text = str(b).strip()
        p.level = 0
    return slide

def _attach_image_to_slide(slide, image_path: str):
    """
    箇条書きスライドの右側に画像を貼る簡易ヘルパー。
    レイアウト次第で多少はみ出す可能性はあるが、
    「とりあえずスライドに紐づける」ことを優先。
    """
    if not image_path:
        return
    try:
        from pptx.util import Inches  # すでに import 済みだが保険
        # 適当な位置・サイズ（右側に縦長で貼るイメージ）
        left = Inches(5.5)
        top = Inches(1.5)
        height = Inches(3.5)
        slide.shapes.add_picture(image_path, left, top, height=height)
    except Exception:
        # 画像が壊れている・パス不正などは無視して続行
        return


def _split_sentences_jp(text: str) -> List[str]:
    # 句点ベースの簡易分割
    text = re.sub(r"\s+", " ", text.strip())
    parts = re.split(r"(。|！|!|？|\?)", text)
    out = []
    buf = ""
    for i in range(0, len(parts), 2):
        s = parts[i].strip()
        punct = parts[i+1] if i+1 < len(parts) else ""
        if not s:
            continue
        out.append((s + punct).strip())
    return out


def _make_outline_from_transcript(text: str) -> List[Dict[str, List[str]]]:
    """
    transcript + OCR テキストから「スライド用アウトライン」を作る。

    戻り値の各要素:
        {
            "title": "スライドタイトル",
            "bullets": ["箇条書き1", "箇条書き2", ...],
        }
    """

    def is_heading(line: str) -> bool:
        """見出しっぽい行かどうかの簡易判定."""
        s = line.strip()
        if not s:
            return False
        # 短めの行を優先（長い文章は本文扱い）
        if len(s) > 30:
            return False

        # よくある見出しキーワード
        keywords = [
            "第", "章", "ステップ", "まとめ", "要約", "ポイント", "ゴール",
            "概要", "目的", "結論", "まとめ", "振り返り",
        ]
        if any(k in s for k in keywords):
            return True

        # 文末が「：」「:」「？」などで終わる短い文も見出し候補
        if s.endswith(("：", ":", "？", "?", "！", "!")):
            return True

        # 英数字だけの短い行（"Agenda", "Intro" など）
        if all(ord(c) < 128 for c in s) and len(s) <= 20:
            return True

        return False

    def split_sentences(block: str) -> List[str]:
        """シンプルな文分割（日本語＋英語混在想定の簡易版）."""
        tmp = (
            block.replace("。", "。\n")
                 .replace("？", "？\n")
                 .replace("！", "！\n")
                 .replace("?", "?\n")
                 .replace("!", "!\n")
        )
        sentences = []
        for line in tmp.splitlines():
            l = line.strip()
            if not l:
                continue
            sentences.append(l)
        return sentences

    # 1) 行単位に分割
    lines = [ln.rstrip() for ln in text.splitlines()]
    # 空行だけのブロックを減らす
    lines = [ln for ln in lines if ln.strip()]

    # 2) 章ごとにまとめる
    sections: List[Dict[str, List[str]]] = []
    current_title: str = ""
    current_body: List[str] = []

    def flush_section():
        nonlocal current_title, current_body
        if not current_title and not current_body:
            return
        # 本文から文を抽出して bullet にする
        bullets: List[str] = []
        for block in current_body:
            for s in split_sentences(block):
                bullets.append(s)
        # 長すぎる bullet は切る
        if len(bullets) > 8:
            bullets = bullets[:8]
        sections.append(
            {
                "title": current_title or (bullets[0] if bullets else "（無題）"),
                "bullets": bullets[1:] if current_title and len(bullets) > 1 else bullets,
            }
        )
        current_title = ""
        current_body = []

    for ln in lines:
        if is_heading(ln):
            # それまでのセクションを確定
            flush_section()
            current_title = ln.strip()
        else:
            current_body.append(ln)

    # 最後のセクション
    flush_section()

    # 3) 章が1つも作れなかった場合は、全体をざっくり分割
    if not sections:
        all_sentences = split_sentences(text)
        if not all_sentences:
            return []

        chunk_size = 6
        for i in range(0, len(all_sentences), chunk_size):
            chunk = all_sentences[i : i + chunk_size]
            title = chunk[0]
            bullets = chunk[1:]
            sections.append(
                {
                    "title": title,
                    "bullets": bullets,
                }
            )

    return sections



def _extract_json_from_text(text: str) -> Optional[dict]:
    """
    LLM応答からJSONだけ取り出す。
    ```json ... ``` や 余分なテキストを含んでも頑張って抜く。
    """
    if not text:
        return None
    # フェンス内json
    m = re.search(r"```json\s*(\{.*?\})\s*```", text, re.S | re.I)
    if m:
        try:
            return json.loads(m.group(1))
        except Exception:
            pass
    # 最初の { から最後の } までを大括弧バランスで抽出
    start = text.find("{")
    end = text.rfind("}")
    if start != -1 and end != -1 and end > start:
        chunk = text[start:end+1]
        try:
            return json.loads(chunk)
        except Exception:
            pass
    # そのままJSONの可能性
    try:
        return json.loads(text)
    except Exception:
        return None



def build_deck(out_path: str,
               title: str,
               transcript: str,
               bullets: Optional[List[List[str]]] = None,
               images: Optional[List[str]] = None,
               use_ai: bool = False) -> str:
    """
    シンプル生成（ローカル）：transcript から素朴にアウトラインを作り、PPTXを保存。
    - out_path: 出力先ファイルパス
    - title: タイトルスライドのタイトル
    - transcript: 元テキスト
    - bullets/images/use_ai: 互換性のための引数
        * bullets: 外部で既に決めたアウトラインがあれば優先（未使用なら None）
        * images: キーフレーム画像のパスリスト（スライドに順番に紐付け）
    戻り値: 保存した PPTX のパス
    """
    prs = _mk_pres()

    # 0) アウトライン作成
    if bullets is not None:
        # 外部から明示的に bullet リストを渡された場合はそれを優先する設計余地
        outline = [
            {"title": f"スライド {i+1}", "bullets": b or ["（内容なし）"]}
            for i, b in enumerate(bullets)
        ]
    else:
        outline = _make_outline_from_transcript(transcript)

    # 箇条書きをスライド用に整形
    for sec in outline:
        sec["bullets"] = _normalize_bullets(sec.get("bullets", []))

    # 1) タイトルスライド
    _add_title_slide(prs, title)

    # 2) アジェンダスライド（最初の数枚のタイトルを並べる）
    agenda_titles = [s["title"] for s in outline[:min(8, len(outline))]]
    if agenda_titles:
        _add_bullet_slide(prs, "アジェンダ", agenda_titles)

    # 3) 本文スライド（必要なら画像を紐付け）
    img_list = images or []
    for idx, s in enumerate(outline):
        slide = _add_bullet_slide(prs, s["title"], s["bullets"])

        if idx < len(img_list):
            _attach_image_to_slide(slide, img_list[idx])

    # 4) まとめスライド（各スライドのタイトルを最後にもう一度並べる）
    if outline:
        summary_bullets = [s["title"] for s in outline[-min(8, len(outline)):]]
        _add_bullet_slide(prs, "まとめ", summary_bullets)

    return _safe_save(prs, out_path)




def build_deck_from_timeline(
    out_path: str,
    title: str,
    entries: List[Dict[str, Any]],
) -> str:
    """
    ASR + OCR + 画像をまとめた timeline entries からスライドを生成する。

    entries: [
      {
        "asr_text": "...この辺で話していた内容...",
        "ocr_text": "...スライドの文字...",
        "image_path": "path/to/frame.png" or None,
      },
      ...
    ]
    """
    prs = _mk_pres()

    # 1) タイムラインから「アウトライン用テキスト」を作る
    text_chunks: List[str] = []
    for e in entries:
        asr = str(e.get("asr_text", "")).strip()
        ocr = str(e.get("ocr_text", "")).strip()

        chunk_lines: List[str] = []
        if asr:
            chunk_lines.append(asr)
        if ocr:
            chunk_lines.append("【スライド内のテキスト】")
            chunk_lines.append(ocr)

        if chunk_lines:
            text_chunks.append("\n".join(chunk_lines))

    if text_chunks:
        outline_source = "\n\n".join(text_chunks)
    else:
        outline_source = "内容が抽出できませんでした。"

    # 2) アウトライン生成（章構造＋見出し検出）
    outline = _make_outline_from_transcript(outline_source)
    
    # 箇条書きをスライド用に整形
    for sec in outline:
        sec["bullets"] = _normalize_bullets(sec.get("bullets", []))


    # 3) タイトルスライド
    _add_title_slide(prs, title)

    # 4) アジェンダスライド（最初の数章のタイトル）
    agenda_titles = [s["title"] for s in outline[: min(8, len(outline))]]
    if agenda_titles:
        _add_bullet_slide(prs, "アジェンダ", agenda_titles)

    # 5) 本文スライド + 画像
    #    - スライドごとに bullets が長すぎる場合は自動的に「(続き)」スライドを作る
    #    - 画像は outline と entries のインデックス比でざっくり割り当てる
    n_entries = len(entries) if entries else 0
    total_sections = len(outline) if outline else 0
    max_bullets_per_slide = 6

    slide_index = 0

    for sec_idx, sec in enumerate(outline):
        title_text = sec["title"]
        bullets = sec.get("bullets", []) or []

        if not bullets:
            # 箇条書きがない場合でも、とりあえずタイトルだけのスライドを出す
            bullets_for_slide = []
            slide = _add_bullet_slide(prs, title_text, bullets_for_slide)
            # 画像割り当て
            if n_entries > 0 and total_sections > 0:
                j = int(slide_index * n_entries / max(1, total_sections))
                j = min(j, n_entries - 1)
                img_path = entries[j].get("image_path")
                if img_path:
                    _attach_image_to_slide(slide, img_path)
            slide_index += 1
            continue

        # 箇条書きが多い場合は複数スライドに自動分割
        for offset in range(0, len(bullets), max_bullets_per_slide):
            chunk = bullets[offset : offset + max_bullets_per_slide]
            if offset == 0:
                t = title_text
            else:
                t = f"{title_text}（つづき）"

            slide = _add_bullet_slide(prs, t, chunk)

            # 画像割り当て:
            # スライド通し番号 slide_index を使って entries に対応付ける
            if n_entries > 0 and total_sections > 0:
                j = int(slide_index * n_entries / max(1, total_sections * 2))
                # *2 はざっくり「分割スライドも含めて」調整するための係数
                j = max(0, min(j, n_entries - 1))
                img_path = entries[j].get("image_path")
                if img_path:
                    _attach_image_to_slide(slide, img_path)

            slide_index += 1

    # 6) まとめスライド
    if outline:
        summary_bullets = [s["title"] for s in outline[-min(8, len(outline)) :]]
        _add_bullet_slide(prs, "まとめ", summary_bullets)

    return _safe_save(prs, out_path)




def build_deck_from_transcript_with_gemini(transcript: str,
                                           context: str = "",
                                           out_path: str = "output.pptx",
                                           title: str = "自動作成資料") -> str:
    """
    Gemini を使って transcript からアウトライン(JSON)を作成し、PPTXを出力。
    - 失敗したらローカル版(build_deck)にフォールバック。
    """
    # 1) Gemini が使えない場合はフォールバック
    if get_client is None:
        return build_deck(out_path, title, transcript)

    # 2) プロンプト
    model = os.environ.get("GEMINI_MODEL", "gemini-2.0-flash")
    prompt = (
        "あなたは日本語のプレゼン資料作成アシスタントです。"
        "以下の講義・会議の文字起こしから、受講者に配布するスライド資料を作成します。"
        "スライドはわかりやすく、過剰な文章は避け、要点だけを箇条書きでまとめてください。"
        "目安として 8〜16 枚程度のスライド構成にしてください（とても短い場合は 5 枚程度でもかまいません）。"
        "各スライドの箇条書きは 3〜5 個とし、1つの箇条書きは40文字以内を目安にしてください。"
        "箇条書きは「〜すること」「〜のポイント」のような名詞・体言止めで、文末に句点「。」を付けないでください。"
        "同じ内容を複数のスライドで繰り返さないようにしてください。"
        "スライド構成には、できれば次の要素を含めてください: "
        "1枚目付近に「アジェンダ」スライド、最後に「まとめ」スライド。"
        "必要に応じて「キーメッセージ」や「今後のアクション」をまとめるスライドを追加してもかまいません。"
        "出力は JSON のみとし、マークダウンのコードフェンスや説明文は書かないでください。\n\n"
        "期待するJSONスキーマ:\n"
        "{\n"
        '  "slides": [\n'
        '    {"title": "スライドタイトル", "bullets": ["要点1","要点2","要点3"]},\n'
        "    ...\n"
        "  ]\n"
        "}\n\n"
        "全てのスライドは、この `slides` 配列に順番に入れてください。\n\n"
    )

    if context:
        prompt += f"コンテキスト/用途: {context}\n\n"
    # 長すぎる transcript を切る（API安全策）
    t = transcript.strip()
    if len(t) > 8000:
        t = t[:8000] + "\n（以下省略）"
    prompt += "=== 対象テキスト ===\n" + t

    try:
        client = get_client()
        # google-genai の想定API（.models.generate_content）
        resp = client.models.generate_content(
            model=model,
            contents=prompt
        )
        # 応答テキストの取り出し（実装差異に広く対応）
        text = None
        for attr in ("text", "output_text"):
            if hasattr(resp, attr):
                text = getattr(resp, attr)
                if text:
                    break
        if text is None:
            # candidates 経由
            try:
                cand = resp.candidates[0]
                parts = getattr(cand, "content", cand).parts
                text = "\n".join([getattr(p, "text", "") for p in parts if getattr(p, "text", "")])
            except Exception:
                text = str(resp)

        data = _extract_json_from_text(text)
        if not data or "slides" not in data or not isinstance(data["slides"], list):
            # 構造が取れなければローカル生成にフォールバック
            return build_deck(out_path, title, transcript)

        # 3) PPTX 作成
        prs = _mk_pres()
        _add_title_slide(prs, title, subtitle="(Gemini 生成)")

        for s in data["slides"]:
            st = str(s.get("title", ""))
            bl = s.get("bullets", [])
            bl = [str(b) for b in bl if str(b).strip()]
            # 箇条書きをスライド用に整形
            bl = _normalize_bullets(bl)

            if not st:
                st = "スライド"
            if not bl:
                bl = ["（内容なし）"]
            _add_bullet_slide(prs, st, bl)


        return _safe_save(prs, out_path)

    except Exception:
        # 何かあっても落とさない
        return build_deck(out_path, title, transcript)
