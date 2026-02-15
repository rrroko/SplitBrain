# app/core/doc_ingest.py
from __future__ import annotations

from pathlib import Path
from typing import List

from pptx import Presentation

try:  # PDF用
    import pdfplumber  # pip install pdfplumber
except Exception:  # ランタイムで「入ってないよ」と伝える
    pdfplumber = None


# -------------------------------------------------------------
# PDF
# -------------------------------------------------------------
def extract_pdf_pages(pdf_path: str | Path) -> List[str]:
    """
    PDFファイルを開いて、ページごとのテキストを返す。

    失敗したページは空文字になる。
    pdfplumber がインストールされていない場合は RuntimeError を投げる。
    """
    if pdfplumber is None:
        raise RuntimeError(
            "pdfplumber がインストールされていません。"
            "  次のコマンドでインストールしてください:\n"
            "  python -m pip install pdfplumber"
        )

    path = Path(pdf_path)
    texts: List[str] = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            try:
                txt = page.extract_text() or ""
            except Exception:
                txt = ""
            texts.append(txt)
    return texts


# -------------------------------------------------------------
# PPTX
# -------------------------------------------------------------
def extract_pptx_slides(pptx_path: str | Path) -> List[str]:
    """
    既存PPTXからスライドごとのテキストを取り出す。
    - テキストボックス
    - 表セルのテキスト
    をつなげて1スライド分の文字列にする。
    """
    path = Path(pptx_path)
    prs = Presentation(path)
    slides_text: List[str] = []

    for slide in prs.slides:
        parts: List[str] = []
        for shape in slide.shapes:
            try:
                if hasattr(shape, "text") and shape.text:
                    parts.append(shape.text)
                elif getattr(shape, "has_table", False):
                    tbl = shape.table
                    for row in tbl.rows:
                        for cell in row.cells:
                            if cell.text:
                                parts.append(cell.text)
            except Exception:
                continue
        slides_text.append("\n".join(p for p in parts if p))
    return slides_text


# -------------------------------------------------------------
# PDF 内の埋め込み画像を個別ファイルとして取り出す
# -------------------------------------------------------------
try:
    import fitz  # PyMuPDF
except Exception:
    fitz = None


def extract_pdf_inline_images(pdf_path: str | Path, out_dir: str | Path) -> list[str]:
    """
    PDF内の埋め込み画像だけを out_dir 以下に保存し、
    生成したファイルパスのリストを返す。

    例: outputs/pdf_images/xxx_p1_img1.png, xxx_p1_img2.jpg ...
    """
    if fitz is None:
        raise RuntimeError(
            "PyMuPDF(fitz) がインストールされていません。"
            "  python -m pip install pymupdf を実行してください。"
        )

    pdf_path = Path(pdf_path)
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)

    doc = fitz.open(str(pdf_path))
    stem = pdf_path.stem
    paths: list[str] = []

    for page_index in range(len(doc)):
        page = doc[page_index]
        image_list = page.get_images(full=True)
        if not image_list:
            continue

        for img_idx, img in enumerate(image_list, start=1):
            xref = img[0]
            base_image = doc.extract_image(xref)
            img_bytes = base_image["image"]
            ext = base_image.get("ext", "png")
            out_file = out_dir / f"{stem}_p{page_index+1}_img{img_idx}.{ext}"
            with open(out_file, "wb") as f:
                f.write(img_bytes)
            paths.append(str(out_file))

    return paths
